import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches


def generate_cold_style_ppt(raw_content_string, output_filename="Cold_Logic_Presentation_V2.pptx"):
    """
    将原始文本内容转换为黑底白字风格的PPT。
    每行文本将成为一张独立幻灯片的标题，并实现居中和字体控制。
    """
    # 1. 解析内容：将多行字符串拆分为列表
    slide_contents = [line.strip() for line in raw_content_string.split('\n') if line.strip()]

    if not slide_contents:
        print("错误：未检测到有效内容行，请确保输入内容不为空。")
        return

    # 2. 初始化PPT
    prs = Presentation()

    # --- 关键的冷峻风格设置 ---
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)

    # 您可以更换为您本地已安装的任何冷峻风格字体（如 Arial Narrow, Consolas, 或其他极简字体）
    TARGET_FONT_NAME = 'Segoe Boot Semilight'
    TARGET_FONT_SIZE = Pt(80)

    # 3. 使用空白布局，以实现最大化的居中掌控
    blank_slide_layout = prs.slide_layouts[6]

    print(f"开始处理 {len(slide_contents)} 行内容...")

    # 4. 循环生成幻灯片
    for i, content in enumerate(slide_contents):
        slide = prs.slides.add_slide(blank_slide_layout)

        # 设置幻灯片背景为纯黑色
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BLACK

        # 创建一个覆盖大部分幻灯片的文本框，以便实现全局居中
        # 居中位置：Left/Top/Width/Height
        left = Inches(0.5)
        top = Inches(0.5)
        width = prs.slide_width - Inches(1.0)
        height = prs.slide_height - Inches(1.0)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        # --- 居中控制 ---
        # 1) 文本框内文字垂直居中
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 2) 文本框大小适应文字，但我们仍使用大框来确保视觉居中
        # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # (这里不使用自动适应，以保持文本框的固定大架构)

        # 插入内容并设置风格
        p = tf.paragraphs[0]
        p.text = content

        # --- 字体控制 ---
        font = p.font
        font.size = TARGET_FONT_SIZE
        font.color.rgb = WHITE
        font.name = TARGET_FONT_NAME

        # 可选：移除段落之间的默认间距，增加极简感
        p.space_after = Pt(0)

        # 5. 保存文件
    prs.save(output_filename)
    print(f"\n✅ 转换完成！PPT文件已保存为: {output_filename}")
    print(f"应用风格：{TARGET_FONT_NAME}, 字号 {int(TARGET_FONT_SIZE.pt)} 磅，黑底白字，完全居中。")


if __name__ == "__main__":
    # --- 您的 Word 内容请替换到这里 ---
    word_content = """

















































































































































































































































































































































































































































































































































































































"""

    # 调用函数生成PPT
    generate_cold_style_ppt(word_content)